from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def read_text_file(file_path):
    data = {}
    with open(file_path, 'r', encoding='utf-8') as file:
        lines = file.readlines()[1:]  # 跳过标题行
        for line in lines:
            category, num, content = line.strip().split('\t')
            if category not in data:
                data[category] = []
            data[category].append((num, content))
    # 数据校验
    for category, items in data.items():
        if len(items) != 25:
            raise ValueError(f"类别 '{category}' 的数据条数应为 25，实际为 {len(items)}")
    return data

def create_ppt(data):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 默认宽
    prs.slide_height = Inches(7.5)   # 默认高
    slide_layout = prs.slide_layouts[6]  # 空白幻灯片

    for category, items in data.items():
        slide = prs.slides.add_slide(slide_layout)
        # 设置背景颜色
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x26, 0x26, 0x26)

        # 添加类别标题
        title = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(11.33), Inches(0.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = category
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = 1

        # 调整圆形布局参数
        circle_radius = Inches(0.6)  # 减小圆形大小
        horizontal_spacing = Inches(0.2)
        vertical_spacing = Inches(0.2)
        left_margin = Inches(1.5)
        top_margin = Inches(1)

        for i in range(5):
            for j in range(5):
                index = i * 5 + j
                num, content = items[index]
                text = f"{num}. {content}"

                # 计算位置
                left = left_margin + j * (circle_radius * 2 + horizontal_spacing)
                top = top_margin + i * (circle_radius * 2 + vertical_spacing)

                # 绘制圆形
                shape = slide.shapes.add_shape(
                    9,  # 圆形
                    left, top,
                    circle_radius * 2, circle_radius * 2
                )
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0x48, 0x74, 0xcb)

                # 添加文本
                tf = shape.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = Pt(12)  # 调整字体大小
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = 1  # 居中

    prs.save('output.pptx')

if __name__ == "__main__":
    file_path = 'ppt草稿.txt'
    data = read_text_file(file_path)
    create_ppt(data)